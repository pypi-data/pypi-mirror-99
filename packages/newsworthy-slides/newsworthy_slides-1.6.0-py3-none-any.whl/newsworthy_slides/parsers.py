import os.path
import requests
import tempfile
import re
from datetime import datetime
from bs4 import (BeautifulSoup,
                 NavigableString,
                 Tag)
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER 
from pptx.enum.text import MSO_AUTO_SIZE
from pptx import Presentation
from pptx.slide import Slide, NotesSlide

from newsworthy_slides.utils import (move_slide, inner_html, get_placeholder,
                                     get_slide_layout)
from newsworthy_slides.validationerror import (PlaceholderContentInvalid,
                                               PlaceholderTypeInvalid,
                                               TooManyPlaceholder,
                                               SlideLayoutInvalid,
                                               SlidePositionInvalid,
                                               NoSuchImage,
                                               InvalidMetadata)

def add_slide(slide_xml, pres):
    """Takes in the xml for a slide and a new slide to add
    to a base presentation.

    Parses all placeholders in the xml slide and creates
    shapes for the new slide.

    :param slide_xml: a html string or a BeautifulSoup element with `slide` has
        the root tag.
    :returns: pptx.Slide object
    """
    if isinstance(slide_xml, Tag):
        slide_soup = slide_xml
    else:
        slide_soup = BeautifulSoup(slide_xml, 'html.parser')

    if slide_soup.name != "slide":
        slide_soup = slide_soup.select_one("slide")
        if slide_soup is None:
            raise Exception(f"No slide tag found in\n {slide_xml}")

    # Set up and add the new slide
    slide_layout_name = slide_soup.get("layout")
    if slide_layout_name is None:
        raise SlideLayoutInvalid(f"Slide layout attribute missing in\n{slide_soup}")

    slide_layout = get_slide_layout(slide_layout_name, pres)
    slide_obj = pres.slides.add_slide(slide_layout)

    # Parse placeholders
    ph_list = slide_soup.select("placeholder")
    for i, placeholder_soup in enumerate(ph_list):
        placeholder_name = placeholder_soup.get("name")

        if placeholder_name is None:
            # If no placeholder name is set, get placeholders in order instead.
            try:
                placeholder_obj = list(slide_obj.placeholders)[i]
            except IndexError:
                msg = (f"Unable to add placeholder number {i+1}. There are only "
                       f"{len(slide_obj.placeholders)} placeholders in the slide"
                       f" layout {slide_obj.slide_layout.name}.\n\n"
                       f"{slide_soup}.")
                raise TooManyPlaceholder(msg)
        else:
            placeholder_obj = get_placeholder(placeholder_name, slide_obj)

        ph_type = placeholder_soup.get("type")
        if ph_type == "text":
            new_ph = _parse_text_placeholder(placeholder_soup, placeholder_obj)
        elif ph_type == "image":
            new_ph = _parse_image_placeholder(placeholder_soup, placeholder_obj, slide_obj)
        elif ph_type == "table":
            new_ph = _parse_table_placeholder(placeholder_soup, placeholder_obj)
        elif ph_type is None:
            msg = f"Placeholder type attribute missing in\n {placeholder_soup}"
            raise PlaceholderTypeInvalid(ph_type)
        else:
            msg = f"Unable to parse {ph_type} placeholders."
            raise NotImplementedError(msg)

    # Determine position of new slide
    # python-pptx does not (yet) have built in way to easily postion new slides
    # Instead we first add the slide to the end (done above), and then move it
    position = slide_soup.get("position")
    if position is not None:
        position = int(position)
        n_slides = len(pres.slides)
        if position > n_slides:
            msg = f"There are only {n_slides} in presentation:\n{slide_soup}"
            raise SlidePositionInvalid(msg)
        move_slide(pres, -1, position)

    # Parse notes
    notes_soup = slide_soup.select_one("notes")
    if notes_soup:
        notes_obj = _parse_notes(notes_soup, slide_obj)

    return slide_obj

def add_metadata(presentation_xml, presentation):
    """Parses presentation level metadata such as title an and author from
    a presentation tag.

    Available meta data is:
        author: string – An entity primarily responsible for making the content
            of the resource.
        category: string – A categorization of the content of this package.
            Example values might include: Resume, Letter, Financial Forecast, Proposal, or Technical Presentation.
        comments: string – An account of the content of the resource.
        content_status: string – completion status of the document, e.g. ‘draft’
        created: datetime – time of intial creation of the document
        identifier: string – An unambiguous reference to the resource
            within a given context, e.g. ISBN.
        keywords: string – descriptive words or short phrases likely to be
            used as search terms for this document
        language: string – language the document is written in
        last_modified_by: string – name or other identifier (such as email
            address) of person who last modified the document
        last_printed: datetime – time the document was last printed
        modified: datetime – time the document was last modified
        revision: int – number of this revision, incremented by the
            PowerPoint® client once each time the document is saved. Note
            however that the revision number is not automatically incremented
            by python-pptx.
        subject: string – The topic of the content of the resource.
        title: string – The name given to the resource.
        version: string – free-form version string
    (from https://python-pptx.readthedocs.io/en/latest/api/presentation.html#coreproperties-objects)
    """
    core_props = presentation.core_properties
    p = presentation
    allowed_properties = [
        ("author", str),
        ("category", str),
        ("comments", str),
        ("content_status", str),
        ("created", datetime),
        ("identifier", str),
        ("keywords", str),
        ("language", str),
        ("last_modified_by", str),
        ("last_printed", datetime),
        ("modified", datetime),
        ("revision", int),
        ("subject", str),
        ("title", str),
        ("version", str),
    ]


    if isinstance(presentation_xml, Tag):
        pres_soup = presentation_xml
    else:
        pres_soup = BeautifulSoup(presentation_xml, 'html.parser')

    metadata_tag = pres_soup.select_one("presentation")
    if metadata_tag:
        for prop_name, prop_type in allowed_properties:
            if prop_name in metadata_tag.attrs:
                value = metadata_tag[prop_name]

                if prop_type.__name__ == "datetime":
                    # TODO: handle hours and minutes
                    try:
                        value = datetime.strptime(value, "%Y-%m-%d")
                    except ValueError:
                        msg = f"Unable to parse datetime from {value}"
                        raise InvalidMetadata(msg)
                else:
                    try:
                        value = prop_type(value)
                    except ValueError:
                        msg = f"Unable to parse {prop_type} from {value}"
                        raise InvalidMetadata(msg)

                setattr(p.core_properties, prop_name, value)

    return presentation




def _parse_text_placeholder(placeholder_soup, slide_shape):
    """Parse the content of a text placeholder.

    Should return "something" that a text frame can digest.
    https://python-pptx.readthedocs.io/en/latest/user/text.html
    """
    ## Set slide text frame
    if isinstance(slide_shape, NotesSlide):
        text_frame = slide_shape.notes_text_frame
    else:
        text_frame = slide_shape.text_frame
    auto_size = False

    if placeholder_soup.name == "placeholder":
        placeholder = placeholder_soup
    else:
        placeholder = placeholder_soup.select_one("placeholder")

    if placeholder is not None:
        placeholder_content = inner_html(placeholder)
    else:
        placeholder_content = placeholder
    
    if "auto-size" in placeholder.attrs:
        auto_size_opts = {
            "text-to-fit-shape": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
            "shape-to-fit-text": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        }
        try:
            auto_size = auto_size_opts[placeholder.attrs["auto-size"]]
        except KeyError:
            raise ValueError(f"Attribute 'auto-size' must be one of: {list(auto_size_opts.keys())}")

    # p and li tags are interpreted as paragraphs.
    # bullet point design should happen in base presentation
    content_soup = BeautifulSoup(placeholder_content, "html.parser")
    paragraph_list = content_soup.select("p,li")

    # If no p or li tags are presented, the placeholder content is treated
    # as one single paragraph
    if len(paragraph_list) == 0:
        paragraph_list = [content_soup]

    for i, paragraph_soup in enumerate(paragraph_list):
        if i == 0:
            # There is always one empty parapgrah in the slide
            text_frame.clear()
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p_content = inner_html(paragraph_soup).strip()
        _add_runs_to_paragraph(p, p_content)
    
    if auto_size:
        text_frame.auto_size = auto_size

    return slide_shape


def _parse_image_placeholder(placeholder_soup, placeholder_obj, slide_obj):
    """Adds an image to an image or object placeholder. 
    The placeholder type has to be set manually in Powerpoint. 

    To be able to handle resizing and insertion of images to general purpose
    placeholders the shape will be _replaced_ (rather than populated)
    by the image. 

    Supports both online and local image sources.

    Images may be rotated by passing rotation as inline css. For example:
    <img src="an_arrow.png" style="transform:rotate(45);">


    :param placeholder_soup: html containing an img tag.
    :param placeholder_obj: A picture or object placeholder.
    :param slide_obj (pptx.Slide): The slide is needed only to for our size fit
        hack (see code below). Once python-pptx supports auto resizing of
        picture placeholders this argument may be removed.
    """
    placeholder_type = placeholder_obj.placeholder_format.type
    allowed_types = [
        PP_PLACEHOLDER.OBJECT, # general purpose placeholder
        PP_PLACEHOLDER.PICTURE, # image placeholder
    ]
    if placeholder_type not in allowed_types:
        msg = (f"You are trying to add an image to an {placeholder_type} "
               f"placeholder (named {placeholder_obj.name}). "
               f"You will have to modify the slide layout by hand in Powerpoint."
               )
        raise PlaceholderContentInvalid(msg)
    
    # default alignment
    placement = {
        "va": "center",
        "ha": "center",
    }
    for attr in ["va", "vertical-alignment"]:
        if attr in placeholder_soup.attrs:
            placement["va"] = placeholder_soup.attrs[attr]

    for attr in ["ha", "horizontal-alignment"]:
        if attr in placeholder_soup.attrs:
            placement["ha"] = placeholder_soup.attrs[attr]

    
    img_soup = placeholder_soup.select_one("img")
    img_src = img_soup["src"]
    # If you have an image url
    if img_src.startswith("http") or img_src.startswith("www"):
        r = requests.get(img_src, stream=True)
        try:
            r.raise_for_status()
        except:
            raise NoSuchImage(f"Unable to download image from {img_src}")

        r.raw.decode_content = True
        with tempfile.TemporaryDirectory() as tmppath:
            filepath = os.path.join(tmppath, "image")
            with open(filepath, "wb") as f:
                for chunk in r.iter_content(chunk_size=128):
                    f.write(chunk)

            # Hack: .insert_picture() does not (yet) support scaling of images inserted
            # to image placeholders. This is workaound replaces the placeholder with
            # a properly fitted image shape in the same place.
            # Possible patch coming up: https://github.com/scanny/python-pptx/pull/439
            added_pic = _replace_and_fit(filepath, placeholder_obj, slide_obj, **placement)

    # If you have a local image file
    else:
        if not os.path.exists(img_src):
            raise NoSuchImage(img_src)

        added_pic = _replace_and_fit(img_src, placeholder_obj, slide_obj, **placement)

    # Apply rotation (if available in inline css)
    # TODO: Make other shapes rotateable
    if "style" in img_soup.attrs:
        import cssutils
        styles = cssutils.parseStyle(img_soup["style"])
        if "transform" in styles and "rotate" in styles["transform"]:
            # TODO: Not sure this covers all allowed css syntaxes for rotate()
            rot_degrees = re.search("rotate\((-?\d+)\)", styles["transform"]).group(1)
            added_pic.rotation = float(rot_degrees)

    return placeholder_obj

def _parse_table_placeholder(placeholder_soup, placeholder_obj):
    """Parse the content of an html table and adds a PPT Table. Note that the
    placeholde explicitly has to be a _table placeholder_. This has to be set
    manually in Powerpoint. At the moment of writing Google Slides does not
    support table placeholders. A base presentation from Google Slides will
    not, in other words, be able to handle this funtion.

    Cells may contain some html. The content is parsed like text paragraphs.

    :param placeholder_soup: html containing a table tag.
    :param placeholder_obj: A table placeholder.
    """

    # validate placeholder
    placeholder_type = placeholder_obj.placeholder_format.type
    if placeholder_type != 12:
        msg = (f"You are trying to add a table to an {placeholder_type} "
               f"placeholder (named {placeholder_obj.name}). "
               f"You will have to modify the slide layout by hand in Powerpoint."
               )
        raise PlaceholderContentInvalid(msg)

    # setup
    table_soup = placeholder_soup.select_one("table")
    rows_soup = table_soup.select("tr")
    header_soup = rows_soup[0].find_all(["th","td"])

    n_rows = len(rows_soup)
    n_cols = len(header_soup)

    shape = placeholder_obj.insert_table(rows=n_rows, cols=n_cols)
    table_obj = shape.table

    # start adding content
    for i, row in enumerate(rows_soup):
        for j, cell_soup in enumerate(row.find_all(["th","td"])):
            cell_obj = table_obj.cell(i,j)

            cell_content = inner_html(cell_soup)
            text_frame = cell_obj.text_frame

            p = text_frame.paragraphs[0]
            # Right-align numbers
            if "value" in cell_soup.get("class", ""):
                p.alignment = PP_ALIGN.RIGHT
            _add_runs_to_paragraph(p, cell_content)

    return placeholder_obj

def _parse_notes(notes_soup, slide_obj):
    """Adds notes based on content from a <note> tag.
    """
    notes_slide = slide_obj.notes_slide
    # Notes are really slides with similar features as regular slides.
    # Hence we can use the same parser as we do for regular slide placeholder.
    # However, for now we don't support images, tables etc in notes. Instead we assume
    # the notes to only contain text.
    # Hackishly we place it in a <placeholder> tag to be able to re-use  _parse_text_placeholder()
    notes_soup = BeautifulSoup(f'<placeholder type="text">{notes_soup}</placeholder>', 'html.parser')
    _parse_text_placeholder(notes_soup, notes_slide)
    return notes_slide


###
# Utility functions
###
def _add_runs_to_paragraph(p, p_content):
    """Adds and formats "runs" to a paraggraph from a given html string.

    Handles basic text formating (bold, italic, links).
    Does NOT support nested tag structures (ie list of lists)

    :param p (pptx.text.text._Paragraph):
    :param p_content: an html string
    """
    p_content = _stringify(p_content)

    # Hack: <br/> will later be translated to \n linebreaks
    # To avoid tag split at p_soup.children below we make a temporary
    # linebreak mark in the string
    p_content = p_content.replace("<br/>", "[LINEBREAK]")

    p_soup = BeautifulSoup(p_content, "html.parser")

    # .children takes care of text outside tags as well
    parts = p_soup.children

    for text_soup in parts:
        run = p.add_run() # create a text run
        is_tag = bool(text_soup.name)
        if is_tag:
            content = inner_html(text_soup)
        else:
            content = text_soup

        content = _stringify(content)
        run.text = _clean_text(content)
        font = run.font
        if text_soup.name == "strong":
            font.bold = True
        if text_soup.name == "em" or text_soup.name == "i":
            font.italic = True
        if text_soup.name == "a":
            run.hyperlink.address = text_soup.attrs["href"]

    return p

def _stringify(p):
    "Get a string, NavigableString or Tag object and returns a string"
    if isinstance(p, str):
        content = p
    if isinstance(p, bytes):
        content = str(p.decode("utf-8"))
    if isinstance(p, NavigableString):
        content = str(p)
    if isinstance(p, Tag):
        content = p.text
    return content


def _clean_text(text):
    "Removes tags and extra spaces from the text"
    text = text.replace("\n", " ")

    text = text.replace("[LINEBREAK]", "\n")

    if "\n" in text:
        text = "\n".join([t.strip() for t in text.split("\n")])

    return text


def _replace_and_fit(img, shape, slide, ha="center", va="center",):
    """Swap an image or content placeholder with a new image to mimic a width fit behaviour.

    Borrowed from https://github.com/scanny/python-pptx/issues/176#issuecomment-366776162
    """
    pic = slide.shapes.add_picture(img, shape.left, shape.top)

    #calculate max width/height for target size
    ratio = min(shape.width / float(pic.width), shape.height / float(pic.height))

    pic.height = int(pic.height * ratio)
    pic.width = int(pic.width * ratio)

    if ha == "center":
        pic.left = int(shape.left + ((shape.width - pic.width) / 2))
    elif ha == "left":
        pic.left = int(shape.left)
    elif ha == "right":
        pic.left = int(shape.left + (shape.width - pic.width))    
    else:
        raise ValueError(f"Invalid horizontal alignment: {ha}."
                         " Expecting center|left|right")

    if va == "center":
        pic.top = int(shape.top + ((shape.height - pic.height) / 2))
    
    elif va == "top": 
        pic.top = int(shape.top)

    elif va == "bottom": 
        pic.top = int(shape.top + (shape.height - pic.height))

    else:
        raise ValueError(f"Invalid vertical alignment: {va}."
                         " Expecting center|top|bottom")


    placeholder = shape.element
    placeholder.getparent().remove(placeholder)

    return pic
