from bs4 import BeautifulSoup
import pytest
from copy import deepcopy
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.shapetree import (SlideShapes, SlidePlaceholder,
                                   PicturePlaceholder, TablePlaceholder)
from pptx.slide import Slide
from newsworthy_slides.utils import get_slide_layout, get_placeholder
from newsworthy_slides.validationerror import (PlaceholderContentInvalid,
                                               TooManyPlaceholder,
                                               PlaceholderTypeInvalid,
                                               SlideLayoutInvalid,
                                               InvalidMetadata,
                                               NoSuchImage)
from newsworthy_slides.parsers import (add_slide, add_metadata,
                                       _parse_text_placeholder,
                                       _parse_image_placeholder,
                                       _parse_table_placeholder)

BASE_PRESENTATION = "test/data/base_presentation.pptx"
test_pres = Presentation(BASE_PRESENTATION)
title_layout = test_pres.slide_layouts[0]

def test_add_slide():
    slides_html = """
<slide master="Dark background slide" layout="Rubrik och innehåll">
    <placeholder type="text">Schools in Åre face large challenges</placeholder>
    <placeholder type="text">
        <p>Teacher certification drops for a <strong>third year in a row</strong> in Åre.</p>
    </placeholder>
</slide>
    """
    n_slides_before = len(test_pres.slides)

    soup = BeautifulSoup(slides_html, 'xml')
    slide = add_slide(soup, test_pres)
    assert(isinstance(slide, Slide) is True)

    assert(len(slide.placeholders) == 2)

    ph1 = slide.placeholders[0]
    assert(len(ph1.text_frame.paragraphs) == 1)
    p = ph1.text_frame.paragraphs[0]
    assert(p.text=="Schools in Åre face large challenges")


    ph2 = slide.placeholders[1]
    assert(len(ph2.text_frame.paragraphs) == 1)
    p = ph2.text_frame.paragraphs[0]
    assert(len(p.runs) == 3)
    assert(p.runs[0].text.endswith("drops for a "))
    assert(p.runs[2].text.startswith(" in Åre"))

    # make sure a slide was added
    n_slides_after = len(test_pres.slides)
    assert(n_slides_after == n_slides_before + 1)

def test_parse_slide_with_too_many_placeholders():
    html = """
    <slide layout="Endast rubrik">
        <placeholder type="text">1</placeholder>
        <placeholder type="text">2</placeholder>
        <placeholder type="text">3</placeholder>
        <placeholder type="text">4</placeholder>
    </slide>
    """
    with pytest.raises(TooManyPlaceholder):
        slide = add_slide(html, test_pres)


def test_placing_placeholder_content_by_name_and_index():
    # Add placeholders with explicit names
    html = """
    <slide layout="Avsnittsrubrik">
        <placeholder type="text" name="Text Placeholder 2">My points</placeholder>
        <placeholder type="text" name="Title 1">My title</placeholder>
    <slide>
    """
    slide = add_slide(html, test_pres)
    assert(slide.placeholders[0].text=="My title")
    assert(slide.placeholders[1].text=="My points")

    # Add placeholders by index
    html = """
    <slide layout="Avsnittsrubrik">
        <placeholder type="text">My points</placeholder>
        <placeholder type="text">My title</placeholder>
    <slide>
    """
    slide = add_slide(html, test_pres)
    assert(slide.placeholders[1].text=="My title")
    assert(slide.placeholders[0].text=="My points")

def test_add_slide_with_undefined_placeholder():
    html = """
    <slide layout="Avsnittsrubrik">
        <placeholder>My points</placeholder>
    <slide>
    """
    with pytest.raises(PlaceholderTypeInvalid):
        add_slide(html, test_pres)

def test_add_slide_without_layout():
    html = """
    <slide>
        <placeholder>My points</placeholder>
    <slide>
    """
    with pytest.raises(SlideLayoutInvalid):
        add_slide(html, test_pres)

def test_add_slide_with_notes():
    # Add placeholders with explicit names
    html = """
    <slide layout="Avsnittsrubrik">
        <placeholder type="text" name="Title 1">My title</placeholder>
        <notes>
            <li>An important note to keep in mind<li>
        </notes>
    <slide>
    """
    slide = add_slide(html, test_pres)
    assert(slide.has_notes_slide)
    notes_slide = slide.notes_slide
    assert(notes_slide.notes_text_frame.text.strip()=="An important note to keep in mind")

def test_parse_text_with_no_tags():
    html = """
        <placeholder type="text">Hello world</placeholder>
    """
    prs = Presentation()
    test_slide = prs.slides.add_slide(prs.slide_layouts[0])
    shape = test_slide.shapes[0]
    soup = BeautifulSoup(html, 'xml')
    text_ph = _parse_text_placeholder(soup, shape)
    assert(isinstance(text_ph, SlidePlaceholder) is True)

    paragraphs  = text_ph.text_frame.paragraphs
    assert(len(paragraphs) == 1)

    p = paragraphs[0]
    assert(p.text == "Hello world")


def test_parse_text_with_line_breaks():
    html = """
    <placeholder type="text">
      <p>
        Hello<br>
        World!
      </p>
    </placeholder>
    """
    prs = Presentation()
    test_slide = prs.slides.add_slide(prs.slide_layouts[0])
    shape = test_slide.shapes[0]
    soup = BeautifulSoup(html, 'xml')
    text_ph = _parse_text_placeholder(soup, shape)
    assert(isinstance(text_ph, SlidePlaceholder)) is True

    paragraphs  = text_ph.text_frame.paragraphs
    assert(len(paragraphs) == 1)

    assert(paragraphs[0].text) == "Hello\nWorld!"


def test_parse_text_with_p_tags():
    html = """
        <placeholder type="text">
          <p>Point 1</p>
          <p>Point 2</p>
          <p>Point 3</p>
        </placeholder>
    """
    prs = Presentation()
    test_slide = prs.slides.add_slide(prs.slide_layouts[0])
    shape = test_slide.shapes[0]
    soup = BeautifulSoup(html, 'xml')
    text_ph = _parse_text_placeholder(soup, shape)
    assert(isinstance(text_ph, SlidePlaceholder)) is True

    paragraphs  = text_ph.text_frame.paragraphs
    assert(len(paragraphs) == 3)

    assert(paragraphs[0].text == "Point 1")
    assert(paragraphs[1].text == "Point 2")
    assert(paragraphs[2].text == "Point 3")


def test_parse_text_with_li_tags():
    html = """
    <placeholder type="text">
        <li>Point 1</li>
        <li>Point 2</li>
        <li>Point 3</li>
    </placeholder>
    """
    prs = Presentation()
    test_slide = prs.slides.add_slide(prs.slide_layouts[0])
    shape = test_slide.shapes[0]
    soup = BeautifulSoup(html, 'xml')
    text_ph = _parse_text_placeholder(soup, shape)
    assert(isinstance(text_ph, SlidePlaceholder)) is True

    paragraphs  = text_ph.text_frame.paragraphs
    assert(len(paragraphs) == 3)

    assert(paragraphs[0].text == "Point 1")
    assert(paragraphs[1].text == "Point 2")
    assert(paragraphs[2].text == "Point 3")


    # including <ul>
    html2 = """
    <placeholder type="text">
        <ul>
            <li>Point 1</li>
            <li>Point 2</li>
            <li>Point 3</li>
        </ul>
    </placeholder>
    """
    prs2 = Presentation()
    test_slide2 = prs2.slides.add_slide(prs.slide_layouts[0])
    shape2 = test_slide2.shapes[0]
    soup2 = BeautifulSoup(html, 'xml')
    text_ph2 = _parse_text_placeholder(soup2, shape2)
    assert(isinstance(text_ph2, SlidePlaceholder)) is True

    paragraphs  = text_ph.text_frame.paragraphs
    assert(len(paragraphs) == 3)

    assert(paragraphs[0].text == "Point 1")
    assert(paragraphs[1].text == "Point 2")
    assert(paragraphs[2].text == "Point 3")


def test_parse_text_with_strong_tags():
    html = """
    <placeholder type="text">
        <p>
            Hello <strong>world</strong>!
        </p>
    </placeholder>
    """
    prs = Presentation()
    test_slide = prs.slides.add_slide(prs.slide_layouts[0])
    shape = test_slide.shapes[0]
    soup = BeautifulSoup(html, 'xml')
    text_ph = _parse_text_placeholder(soup, shape)
    assert(isinstance(text_ph, SlidePlaceholder)) is True

    paragraphs  = text_ph.text_frame.paragraphs
    assert(len(paragraphs) == 1)

    runs = text_ph.text_frame.paragraphs[0].runs
    assert(len(runs) == 3)
    assert(runs[0].text == "Hello ")
    assert(runs[1].text == "world")
    assert(runs[2].text == "!")


def test_parse_text_with_a_tags():
    # Case 1: Whole placeholder is link
    html = """
    <placeholder type="text">
        <a href='http://www.newsworthy.se'>Visit our site</a>
    </placeholder>
    """
    prs = Presentation()
    test_slide = prs.slides.add_slide(prs.slide_layouts[0])
    shape = test_slide.shapes[0]
    soup = BeautifulSoup(html, 'xml')
    text_ph = _parse_text_placeholder(soup, shape)
    assert(isinstance(text_ph, SlidePlaceholder)) is True

    paragraphs  = text_ph.text_frame.paragraphs
    assert(len(paragraphs) == 1)

    p = paragraphs[0]
    assert(p.text == "Visit our site")
    assert(p.runs[0].hyperlink.address == 'http://www.newsworthy.se')

    # Case 2: Inline link
    html2 = """
    <placeholder type="text">
        Here is <a href='http://www.newsworthy.se'>a link</a> to our world.
    </placeholder>
    """
    prs2 = Presentation()
    test_slide2 = prs2.slides.add_slide(prs.slide_layouts[0])
    shape2 = test_slide2.shapes[0]
    soup2 = BeautifulSoup(html2, 'xml')
    text_ph2 = _parse_text_placeholder(soup2, shape2)
    assert(isinstance(text_ph2, SlidePlaceholder)) is True
    paragraphs  = text_ph2.text_frame.paragraphs
    assert(len(paragraphs) == 1)

    p = paragraphs[0]
    assert(len(p.runs) == 3)
    assert(p.runs[1].hyperlink.address == 'http://www.newsworthy.se')

def test_parse_long_text_with_auto_size():
    html = """
    <placeholder type="text" auto-size="text-to-fit-shape">
        <li>Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod tempor incididunt.</li>
    </placeholder>    
    """
    prs = Presentation()
    test_slide = prs.slides.add_slide(prs.slide_layouts[3])
    shape = test_slide.shapes[0]
    soup = BeautifulSoup(html, 'xml')
    text_ph = _parse_text_placeholder(soup, shape)
    assert text_ph.text_frame.auto_size == 2
    


def test_parse_text_with_multiple_tags():
    html = """
    <placeholder type="text">
        <p>This is <strong>important</strong>.</p>
        <p>Here is <a href='http://www.newsworthy.se'>a link</a> to our world.</p>
    </placeholder>
    """
    prs = Presentation()
    test_slide = prs.slides.add_slide(prs.slide_layouts[0])
    shape = test_slide.shapes[0]
    soup = BeautifulSoup(html, 'xml')
    text_ph = _parse_text_placeholder(soup, shape)
    assert(isinstance(text_ph, SlidePlaceholder)) is True

    paragraphs  = text_ph.text_frame.paragraphs
    assert(len(paragraphs) == 2)

    p1 = text_ph.text_frame.paragraphs[0]
    assert(len(p1.runs) == 3)
    assert(p1.runs[0].text == "This is ")
    assert(p1.runs[1].text == "important")
    assert(p1.runs[1].font.bold is True)
    assert(p1.runs[2].text == ".")

    p2 = text_ph.text_frame.paragraphs[1]
    assert(len(p2.runs) == 3)
    assert(p2.runs[0].text == "Here is ")
    assert(p2.runs[1].text == "a link")
    assert(p2.runs[1].hyperlink.address == "http://www.newsworthy.se")
    assert(p2.runs[2].text == " to our world.")


def test_parse_image_placeholder():
    ###
    # Setup
    ###
    prs = Presentation()
    layout_with_image = prs.slide_layouts[8]

    test_slide = prs.slides.add_slide(layout_with_image)
    text_ph = test_slide.placeholders[0]
    image_ph = test_slide.placeholders[1]

    assert(text_ph.placeholder_format.type == 1)
    assert(image_ph.placeholder_format.type == 18)

    # It is not be possible to add a non-existing image
    html = """
    <placeholder type="image">
        <img src="https://newsworthy.s3.amazonaws.com/robot-writer/charts/THIS_IMG_DOES_NOT_EXIST.png">
    </placeholder>
    """
    soup = BeautifulSoup(html, 'xml')
    with pytest.raises(NoSuchImage):
        _parse_image_placeholder(soup, image_ph, test_slide)


    # It is not be possible to add an image to a non-image placeholder
    with pytest.raises(PlaceholderContentInvalid):
        _parse_image_placeholder(soup, text_ph, test_slide)


    # ...but this should work:
    html = """
    <placeholder type="image">
        <img src="https://newsworthy.s3.amazonaws.com/robot-writer/charts/3110061949.png">
    </placeholder>
    """
    soup = BeautifulSoup(html, 'xml')
    placeholder_obj = _parse_image_placeholder(soup, image_ph, test_slide)

    assert(isinstance(placeholder_obj, PicturePlaceholder) is True)

    #TODO: Not sure how to validate that the content is actually an image

    # Add image wrapped in a figure tag, this should not affect the result
    html = """
    <placeholder type="image">
        <figure class="figchart">
            <img class="chart" src="https://newsworthy.s3.amazonaws.com/robot-writer/charts/3110061949.png">
        </figure>
    </placeholder>
    """
    # p = _parse_image_placeholder(html)
    # TODO: Make assertions


    # parse empty
    html = """<placeholder type="image"></placeholder>"""
    # TODO: Make assertions
    # p = _parse_image_placeholder(html)

def test_parse_image_to_object_placeholder():
    prs = Presentation()
    layout_with_obj_ph = prs.slide_layouts[4]

    test_slide = prs.slides.add_slide(layout_with_obj_ph)
    obj_ph = test_slide.placeholders[4]
    assert(obj_ph.placeholder_format.type == 7)

    html = """
    <placeholder type="image">
        <img src="https://newsworthy.s3.amazonaws.com/robot-writer/charts/3110061949.png">
    </placeholder>
    """
    soup = BeautifulSoup(html, 'xml')
    placeholder_obj = _parse_image_placeholder(soup, obj_ph, test_slide)

    assert(isinstance(placeholder_obj, SlidePlaceholder) is True)


def test_parse_rotated_image():
    ###
    # Setup
    ###
    prs = Presentation()
    layout_with_image = prs.slide_layouts[8]

    test_slide = prs.slides.add_slide(layout_with_image)
    text_ph = test_slide.placeholders[0]
    image_ph = test_slide.placeholders[1]
    # ...but this should work:
    html = """
    <placeholder type="image">
        <img src="https://newsworthy.s3.amazonaws.com/robot-writer/charts/3110061949.png" style="transform:rotate(45)">
    </placeholder>
    """
    soup = BeautifulSoup(html, 'xml')
    _parse_image_placeholder(soup, image_ph, test_slide)
    rotations = [x.rotation for x in test_slide.shapes]
    last_shape = [x for x in test_slide.shapes][-1]

    assert last_shape.rotation == 45.0



def test_parse_table_placeholder():
    html = """
    <table class="table responsive">
     <thead>
      <tr>
       <th scope="col">Kommun</th>
       <th scope="col">Arbetslöshet i juli</th>
       <th scope="col">Förändring (12 mån)</th>
     </tr>
    </thead>
    <tbody>
     <tr>
      <th scope="row">Karlshamn</th>
      <td class="value" data-title="Arbetslöshet i juli" data-value="0.088">8,8 %</td>
      <td class="value" data-title="Förändring (12 mån)" data-value="0.0020000000000000018">+0,2 % <i>up</i></td>
     </tr>
     <tr>
      <th scope="row">Karlskrona</th>
      <td class="value" data-title="Arbetslöshet i juli" data-value="0.072">7,2 %</td>
      <td class="value" data-title="Förändring (12 mån)" data-value="-0.0050000000000000044">−0,5 % <i>down</i></td>
     </tr>
    </tbody>
   </table>
    """
    soup = BeautifulSoup(html,"html.parser")
    prs = deepcopy(test_pres)
    layout_with_table = get_slide_layout("Rubrik och tabell", prs)

    test_slide = prs.slides.add_slide(layout_with_table)

    text_ph = get_placeholder("Title 1", test_slide)
    table_ph = get_placeholder("Table Placeholder 2", test_slide)

    table_ph = _parse_table_placeholder(soup, table_ph)
    assert(isinstance(table_ph, TablePlaceholder))

    table = test_slide.shapes[-1].table
    assert(table.cell(0,0).text == "Kommun")
    cell_with_tag = table.cell(1,2)
    runs = cell_with_tag.text_frame.paragraphs[0].runs
    assert(len(runs) == 2)
    assert(runs[0].text == "+0,2 % ")
    assert(runs[1].text == "up")


def test_parse_table_placeholder_with_th():
    html = """
<table class="table responsive">
<thead>
<tr>
    <th scope="col">Kommun</th>
    <th scope="col">Placering i länet</th>
    <th scope="col">Placering i riket</th>
</tr>
</thead>
<tbody>
<tr class="highlight">
    <th scope="row">Övertorneå</th>
    <td class="value" data-title="Placering i länet" data-value="1">1<small>/14</small></td>
    <td class="value" data-title="Placering i riket" data-value="1">1/290</td>
</tr>
<tr>
    <th scope="row"><a href="https://www.lararforbundet.se/basta-skolkommun/2019/arvidsjaur" target="_top">Arvidsjaur</a></th>
    <td class="value" data-title="Placering i länet" data-value="2">2<small>/14</small></td>
    <td class="value" data-title="Placering i riket" data-value="4">4/290</td>
</tr>
<tr class="ellipsis">
    <td colspan="3">⋮</td>
</tr>
<tr>
    <th scope="row"><a href="https://www.lararforbundet.se/basta-skolkommun/2019/pajala" target="_top">Pajala</a></th>
    <td class="value" data-title="Placering i länet" data-value="12">12<small>/14</small></td>
    <td class="value" data-title="Placering i riket" data-value="197">197/290</td>
</tr>
</tbody>
</table>
    """
    soup = BeautifulSoup(html,"html.parser")
    prs = deepcopy(test_pres)
    layout_with_table = get_slide_layout("Rubrik och tabell", prs)

    test_slide = prs.slides.add_slide(layout_with_table)

    table_ph = get_placeholder("Table Placeholder 2", test_slide)
    table_ph = _parse_table_placeholder(soup, table_ph)
    assert(isinstance(table_ph, TablePlaceholder))

    table = test_slide.shapes[-1].table
    assert(table.cell(1,0).text == "Övertorneå")
    assert(table.cell(1,1).text == "1/14")
    assert(table.cell(1,2).text == "1/290")

def test_add_metadata():
    prs = deepcopy(test_pres)
    html = """
    <presentation author="Jens" title="Test presentation" created="2020-10-01" category="">
    </presentation>
    """
    add_metadata(html, prs)

    assert prs.core_properties.author == "Jens"
    assert prs.core_properties.title == "Test presentation"
    assert prs.core_properties.created == datetime(2020,10,1)
    prs.save("test.pptx")

    with pytest.raises(InvalidMetadata):
        html = """
        <presentation created="bad date">
        </presentation>
        """
        add_metadata(html, prs)
