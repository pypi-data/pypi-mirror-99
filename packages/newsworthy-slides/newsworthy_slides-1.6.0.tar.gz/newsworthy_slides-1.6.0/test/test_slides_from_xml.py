from newsworthy_slides.main import slides_from_xml
from newsworthy_slides.validationerror import SlideLayoutInvalid, SlidePositionInvalid
from pptx import Presentation
from pptx.presentation import Presentation as PresentationClass
import pytest

BASE_PRESENTATION = u"test/data/base_presentation.pptx"



def test_basic():
    xml = """
<slide layout="Rubrikbild">
    <placeholder type="text">Slide 1</placeholder>
</slide>
<slide layout="Rubrik och innehåll">
    <placeholder type="text">Slide 2</placeholder>
</slide>
<slide layout="Två innehållsdelar">
    <placeholder type="text">Slide 3</placeholder>
    <placeholder type="image" vertical-alignment="bottom">
        <img src="test/data/robot-writer-chart.png">
    </placeholder>

    <notes>
        <li>Read more <a href="https://www.newsworthy.se">on our site</a></li>
    </notes>
</slide>
<slide layout="Två innehållsdelar">
    <placeholder type="text">Mycket text här</placeholder>
    <placeholder type="text" auto-size="text-to-fit-shape">
        <li>Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod tempor incididunt.</li>
        <li> ut labore et dolore magna aliqua. Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris</li>
        <li>nisi ut aliquip ex ea commodo consequat. Duis aute irure dolor in reprehenderit</li>
        <li>in voluptate velit esse cillum dolore eu fugiat nulla pariatur.</li>
        <li>Excepteur sint occaecat cupidatat non proident, sunt in culpa qui officia deserunt</li>
        <li>mollit anim id est laborum.</li>
    </placeholder>
</slide>
    """
    pres = slides_from_xml(xml, BASE_PRESENTATION)
    assert(isinstance(pres, PresentationClass))
    n_slides_after = len(pres.slides)

    # There are 3 slides because an empty first slide is also present
    n_slides_before = 2
    assert(n_slides_after == n_slides_before + 4)
    s1 = pres.slides[2]
    s2 = pres.slides[3]
    assert(s1.slide_layout.name == "Rubrikbild")
    assert(s2.slide_layout.name == "Rubrik och innehåll")

    pres.save("test/rendered/test_basic.pptx")

def test_slide_positioning():
    # defult behaviour: append slides to end
    xml = """
<slide layout="Rubrikbild">
    <placeholder type="text">New slide</placeholder>
</slide>
    """
    pres = slides_from_xml(xml, BASE_PRESENTATION)
    last_slide = pres.slides[-1]
    assert(last_slide.placeholders[0].text == "New slide")

    # define position
    xml = """
<slide layout="Rubrikbild" position="0">
    <placeholder type="text">New slide</placeholder>
</slide>
    """
    pres = slides_from_xml(xml, BASE_PRESENTATION)
    new_slide = pres.slides[0]
    assert(new_slide.placeholders[0].text == "New slide")

    # invalid poistion
    xml = """
<slide layout="Rubrikbild" position="999">
    <placeholder type="text">New slide</placeholder>
</slide>
    """
    with pytest.raises(SlidePositionInvalid):
        pres = slides_from_xml(xml, BASE_PRESENTATION)
