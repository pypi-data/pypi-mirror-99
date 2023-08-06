from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.presentation import Presentation as PresentationClass
from newsworthy_slides.parsers import add_slide, add_metadata

def slides_from_xml(slides_xml, base_presentation):
    """Adds slides to a presentation by parsing our custom xml/html markup
    for defining slide content.

    :param slides_xml (str): slide content as XML.
    :param base_presentation: path to the PPTX file that the slides will be
        added to (or and instance of pptx.Presentation).
    """
    if isinstance(base_presentation, PresentationClass):
        pres = base_presentation
    else:
        pres = Presentation(base_presentation)

    add_metadata(slides_xml, pres)

    slides_soup = BeautifulSoup(slides_xml, 'html.parser').select("slide")

    for i, slide_soup in enumerate(slides_soup):
        slide = add_slide(slide_soup, pres)

    return pres
